import fitz
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from pptx import Presentation

model = SentenceTransformer('all-MiniLM-L6-v2')

# 1. Extract text
def extract_text(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_ppt_text(file):
    prs = Presentation(file)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)


# 2. Chunk text
def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    for i in range(0, len(text), chunk_size - overlap):
        chunks.append(text[i:i + chunk_size])
    return chunks

# 3. Create FAISS index
def create_index(chunks):
    embeddings = model.encode(chunks)
    dimension = embeddings.shape[1]

    index = faiss.IndexFlatL2(dimension)
    index.add(np.array(embeddings))

    return index, chunks

# 4. Search
def retrieve(query, index, chunks, k=3):
    query_embedding = model.encode([query])
    distances, indices = index.search(query_embedding, k)

    results = [chunks[i] for i in indices[0]]
    return results